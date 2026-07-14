"""
deck_graphics.py — engine for "Claude directs, an image model illustrates".

Two ways it gets used:
  - API path:  this script calls an image API (gpt-image-1 / Gemini) directly.
  - MCP path:  Claude gets the image from an MCP connector instead, then calls
               this script only for `insert`.

CLI:
  python deck_graphics.py generate "<subject>" out.png [--provider openai|gemini] [--evaluate] [--style Image_Prompts.md]
  python deck_graphics.py insert  deck.pptx <slide_index> img.png out.pptx [--left 6.3 --top 1.4 --width 3.2]
  python deck_graphics.py render  deck.pptx [--slides 1,3,5] [--out render]
  python deck_graphics.py demo    talk.pptx talk_illustrated.pptx

Style resolution for generate/demo: --style FILE  >  ./Image Prompts.md  >  built-in STYLE_GUIDE.

Deps:  pip install openai google-genai anthropic python-pptx pillow markitdown
Keys:  OPENAI_API_KEY  and/or  GEMINI_API_KEY  and/or  ANTHROPIC_API_KEY (for --evaluate)
"""

from __future__ import annotations
import base64, json
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


# ---------------------------------------------------------------------------
# 1. The shared style guide. Edit "Image Prompts.md" instead if you prefer;
#    the skill loads that and passes it in. Define the look ONCE.
# ---------------------------------------------------------------------------
STYLE_GUIDE = """\
Hand-drawn whiteboard doodle illustration, friendly and playful.
Soft rounded line art in charcoal, gentle marker/highlighter fills.
One clear subject, generous whitespace, no photorealism, no 3D, no gradients.
Palette: deep water blue #1E5A8A, teal #2FA4A8, sand #E8D4A0, one warm accent.
Fully TRANSPARENT background so it drops cleanly onto a slide.
If the subject implies a label, render at most 2-3 short words, spelled correctly.
Cohesive with a set: same line weight and character style across all images.
"""

DEFAULT_PROVIDER = "openai"   # gpt-image-1 & Gemini match this style; FLUX/SD don't


def load_style(style_path=None) -> str:
    """--style FILE beats ./Image Prompts.md beats the built-in STYLE_GUIDE."""
    if style_path:
        return Path(style_path).read_text()
    local = Path("Image Prompts.md")
    return local.read_text() if local.exists() else STYLE_GUIDE


def compose_prompt(subject: str, style: str = STYLE_GUIDE) -> str:
    return f"{subject.strip()}\n\nStyle:\n{style.strip()}"


# ---------------------------------------------------------------------------
# 2. Generation
# ---------------------------------------------------------------------------
def _save_png(b64: str, out_path: Path) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_bytes(base64.b64decode(b64))
    return out_path


def generate_openai(prompt: str, out_path: Path, size: str = "1024x1024") -> Path:
    from openai import OpenAI
    client = OpenAI()
    r = client.images.generate(model="gpt-image-1", prompt=prompt, size=size,
                               quality="high", background="transparent")
    return _save_png(r.data[0].b64_json, out_path)


def generate_gemini(prompt: str, out_path: Path, size: str = "1024x1024") -> Path:
    # NOTE: verify the current Gemini image model name at ai.google.dev.
    from google import genai
    client = genai.Client()
    r = client.models.generate_content(model="gemini-2.5-flash-image", contents=prompt)
    for part in r.candidates[0].content.parts:
        if getattr(part, "inline_data", None):
            return _save_png(base64.b64encode(part.inline_data.data).decode(), out_path)
    raise RuntimeError("Gemini returned no image part")


def generate_graphic(subject: str, out_path, provider: str = DEFAULT_PROVIDER,
                     style: str | None = None) -> Path:
    out_path = Path(out_path)
    prompt = compose_prompt(subject, style if style is not None else load_style())
    return {"openai": generate_openai, "gemini": generate_gemini}[provider](prompt, out_path)


def alpha_report(png_path) -> str:
    """One-line transparency check. A transparent PNG often *looks* like it has a
    dark background in image viewers that composite alpha onto black — this tells
    the evaluator (human or Claude) what's really there."""
    from PIL import Image
    img = Image.open(png_path)
    if "A" not in img.getbands():
        return "WARNING: no alpha channel — background is opaque"
    a = img.getchannel("A")
    data = a.get_flattened_data() if hasattr(a, "get_flattened_data") else a.getdata()
    pct = round(100 * sum(1 for v in data if v < 10) / (img.width * img.height))
    return f"{pct}% transparent (a dark backdrop in a viewer is just alpha compositing)"


# ---------------------------------------------------------------------------
# 3. Evaluate. In Cowork/Claude Code, Claude just VIEWS the PNG and decides.
#    This is the programmatic stand-in for unattended runs.
# ---------------------------------------------------------------------------
def critique_graphic(image_path, requirement: str) -> dict:
    import anthropic
    img = Path(image_path).read_bytes()
    msg = anthropic.Anthropic().messages.create(
        model="claude-sonnet-5", max_tokens=300,
        messages=[{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": "image/png",
             "data": base64.b64encode(img).decode()}},
            {"type": "text", "text":
                f"This doodle should satisfy: '{requirement}'. Judge style match, clarity, "
                f'and spelling of any text. Reply with ONLY JSON: '
                f'{{"pass": bool, "reason": str, "fix_hint": str}}'}]}])
    text = "".join(b.text for b in msg.content if b.type == "text").strip()
    return json.loads(text.removeprefix("```json").removesuffix("```").strip())


def generate_until_good(subject, requirement, out_path, provider=DEFAULT_PROVIDER,
                        max_tries=3, style=None) -> Path:
    subj, path = subject, None
    for attempt in range(1, max_tries + 1):
        path = generate_graphic(subj, out_path, provider, style=style)
        v = critique_graphic(path, requirement)
        print(f"  attempt {attempt}: {'PASS' if v['pass'] else 'redo'} — {v['reason']}")
        if v["pass"]:
            return path
        subj = f"{subject}\n\nRevision note: {v['fix_hint']}"
    return path


# ---------------------------------------------------------------------------
# 4. Insertion — transparent PNG onto an existing deck (python-pptx handles it)
# ---------------------------------------------------------------------------
def add_graphic_to_slide(slide, image_path, left_in=6.3, top_in=1.4, width_in=3.2):
    slide.shapes.add_picture(str(image_path), Inches(left_in), Inches(top_in), width=Inches(width_in))


@dataclass
class GraphicSpec:
    slide_index: int
    subject: str
    requirement: str


def build_deck(in_deck, plan, out_deck, art_dir="art", provider=DEFAULT_PROVIDER,
               evaluate=True, style=None):
    prs = Presentation(in_deck)
    for spec in plan:
        png = Path(art_dir) / f"slide{spec.slide_index}.png"
        print(f"slide {spec.slide_index}: {spec.subject[:50]}...")
        if evaluate:
            generate_until_good(spec.subject, spec.requirement, png, provider, style=style)
        else:
            generate_graphic(spec.subject, png, provider, style=style)
        add_graphic_to_slide(prs.slides[spec.slide_index], png)
    prs.save(out_deck)
    print(f"saved -> {out_deck}")
    return out_deck


# ---------------------------------------------------------------------------
# 5. Visual QA — render slides to PNGs so the placed art can be checked for
#    overlap/overflow. macOS: Quick Look (no installs). Elsewhere: LibreOffice
#    to PDF (view the PDF pages instead).
# ---------------------------------------------------------------------------
def render_slides(deck, out_dir="render", slides=None, size=1600) -> list[Path]:
    import platform, shutil, subprocess, tempfile
    out = Path(out_dir); out.mkdir(parents=True, exist_ok=True)
    total = len(Presentation(deck).slides._sldIdLst)
    idxs = list(range(total)) if slides is None else slides

    if platform.system() == "Darwin" and shutil.which("qlmanage"):
        # Quick Look only renders a pptx's FIRST slide, so for each requested
        # index save a temp copy with that slide moved to the front.
        made = []
        with tempfile.TemporaryDirectory() as td:
            for i in idxs:
                prs = Presentation(deck)
                lst = prs.slides._sldIdLst
                sld = list(lst)[i]
                lst.remove(sld); lst.insert(0, sld)
                tmp = Path(td) / f"slide{i}.pptx"
                prs.save(tmp)
                subprocess.run(["qlmanage", "-t", "-s", str(size), "-o", str(out), str(tmp)],
                               capture_output=True, check=True)
                final = (out / f"slide{i}.pptx.png").replace(out / f"slide{i}.png")
                made.append(final); print(final)
        return made

    if shutil.which("soffice"):
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                        "--outdir", str(out), str(deck)], capture_output=True, check=True)
        pdf = out / (Path(deck).stem + ".pdf")
        print(f"{pdf}  (no per-slide PNGs without Quick Look — view this PDF's pages)")
        return [pdf]

    raise RuntimeError("No renderer available: needs macOS Quick Look or LibreOffice (soffice).")


DEMO_PLAN = [
    GraphicSpec(1, "A cute cartoon river gage station on a riverbank measuring water level, "
                   "friendly sensor with a smiling face, tiny label 'stream gage'.",
                   "reads instantly as a stream gage, label spelled right, on-brand blues"),
    GraphicSpec(3, "A doodle water cycle as a simple loop: cloud, rain, river, arrows back up, "
                   "one happy water-droplet character.", "clear closed loop, playful, uncluttered"),
    GraphicSpec(5, "A friendly robot handing a forecast chart to a hydrologist, AI-assisted "
                   "streamflow prediction.", "warm collaborative feel; robot matches the doodle set"),
]


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def _cli():
    import argparse
    p = argparse.ArgumentParser(description="deck-graphics engine")
    sub = p.add_subparsers(dest="cmd", required=True)

    g = sub.add_parser("generate")
    g.add_argument("subject"); g.add_argument("out")
    g.add_argument("--provider", default=DEFAULT_PROVIDER)
    g.add_argument("--evaluate", action="store_true")
    g.add_argument("--requirement", default="matches the house doodle style, clear, spelled correctly")
    g.add_argument("--style", default=None, help="path to a style-guide file (default: ./Image Prompts.md if present)")

    ins = sub.add_parser("insert")
    ins.add_argument("deck"); ins.add_argument("slide", type=int)
    ins.add_argument("image"); ins.add_argument("out")
    ins.add_argument("--left", type=float, default=6.3)
    ins.add_argument("--top", type=float, default=1.4)
    ins.add_argument("--width", type=float, default=3.2)

    r = sub.add_parser("render")
    r.add_argument("deck")
    r.add_argument("--slides", default=None, help="comma-separated 0-based indices (default: all)")
    r.add_argument("--out", default="render")
    r.add_argument("--size", type=int, default=1600)

    d = sub.add_parser("demo")
    d.add_argument("in_deck"); d.add_argument("out_deck")
    d.add_argument("--style", default=None)

    a = p.parse_args()
    if a.cmd == "generate":
        style = load_style(a.style)
        out = (generate_until_good(a.subject, a.requirement, a.out, a.provider, style=style)
               if a.evaluate else generate_graphic(a.subject, a.out, a.provider, style=style))
        print(f"{out}  [{alpha_report(out)}]")
    elif a.cmd == "insert":
        prs = Presentation(a.deck)
        add_graphic_to_slide(prs.slides[a.slide], a.image, a.left, a.top, a.width)
        prs.save(a.out); print(a.out)
    elif a.cmd == "render":
        idxs = [int(s) for s in a.slides.split(",")] if a.slides else None
        render_slides(a.deck, a.out, idxs, a.size)
    elif a.cmd == "demo":
        build_deck(a.in_deck, DEMO_PLAN, a.out_deck, style=load_style(a.style))


if __name__ == "__main__":
    _cli()
